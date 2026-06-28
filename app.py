import os, shutil, io, re, json
from datetime import datetime
import gradio as gr
import pandas as pd

# ================== LANGCHAIN IMPORTS ==================
from langchain_community.document_loaders import (
    Docx2txtLoader,
    UnstructuredPDFLoader,
    UnstructuredExcelLoader,
    UnstructuredImageLoader,
    TextLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_ollama import OllamaEmbeddings, ChatOllama
from langchain_chroma import Chroma

# ================== FILE HANDLING ==================
from docx import Document
from pptx import Presentation

# ================== KONFIGURASI ==================
CHROMA_PATH = "chroma_db"
OUTPUT_DIR = os.path.join(os.getcwd(), "output")
LOG_FILE = os.path.join(OUTPUT_DIR, "rejected_prompts.log")
os.makedirs(OUTPUT_DIR, exist_ok=True)

MODEL_OPTIONS = {
    "Akurat": {
        "word": "qwen2.5:14b-instruct-q4_K_M",
        "excel": "deepseek-coder-v2:16b-lite-instruct-q4_K_M",
        "ppt": "mistral-nemo:12b",
        "jurnal": "qwen2.5:14b-instruct-q4_K_M",
        "laporan": "qwen2.5:14b-instruct-q4_K_M"
    },
    "Hemat": {
        "word": "qwen2.5:14b-instruct-q4_K_M",
        "excel": "deepseek-coder-v2:16b-lite-instruct-q4_K_M",
        "ppt": "mistral-nemo:12b",
        "jurnal": "qwen2.5:14b-instruct-q4_K_M",
        "laporan": "qwen2.5:14b-instruct-q4_K_M"
    }
}

embeddings = OllamaEmbeddings(model="nomic-embed-text:latest")
db = None
last_result = {"content": "", "type": ""}
doc_context = "" # Simpan konteks global buat edukasi

# ================== KNOWLEDGE BASE UNTUK EDUKASI ==================
TERMINOLOGY_EDUCATION = {
    "gaji": {
        "description": "Gaji pokok adalah upah dasar yang diterima pegawai berdasarkan pangkat/golongan",
        "confused_with": ["tunjangan kinerja", "tunjangan jabatan", "tunjangan fungsional"]
    },
    "tunjangan": {
        "description": "Tunjangan adalah tambahan penghasilan di luar gaji pokok",
        "types": {
            "tunjangan kinerja": "Diberikan berdasarkan capaian kinerja dan reformasi birokrasi",
            "tunjangan jabatan": "Diberikan berdasarkan jabatan yang diduduki",
            "tunjangan istri/suami": "Tunjangan untuk pasangan yang sah",
            "tunjangan anak": "Tunjangan untuk anak yang masih menjadi tanggungan"
        }
    },
    "cuti": {
        "description": "Cuti adalah hak pegawai untuk tidak masuk kerja dalam jangka waktu tertentu",
        "types": {
            "cuti tahunan": "Hak cuti rutin 12 hari kerja per tahun",
            "cuti besar": "Cuti panjang setelah bekerja minimal 6 tahun (3 bulan)",
            "cuti melahirkan": "Cuti untuk ibu melahirkan (3 bulan)",
            "cuti di luar tanggungan negara": "Cuti tanpa gaji untuk kepentingan pribadi"
        }
    },
    "izin": {
        "description": "Izin adalah izin tidak masuk kerja dalam waktu singkat",
        "rules": "Maksimal 4 kali dalam 1 bulan, dengan alasan yang sah"
    }
}

def get_education_context(query, doc_context):
    query_lower = query.lower()
    doc_lower = doc_context.lower()
    education_messages = []

    if 'gaji' in query_lower and 'tunjangan' in doc_lower and 'gaji pokok' not in doc_lower:
        education_messages.append({
            "type": "confusion",
            "title": "📚 PERBEDAAN GAJI DAN TUNJANGAN",
            "content": """
**Dokumen ini membahas TUNJANGAN, bukan GAJI POKOK.**

| Istilah | Penjelasan |
| --- | --- |
| **Gaji Pokok** | Upah dasar berdasarkan pangkat/golongan. Tidak dibahas di dokumen ini. |
| **Tunjangan Kinerja** | Tambahan penghasilan berdasarkan kinerja. Dibahas di dokumen ini. |

💡 **Saran:** Upload SK penggajian jika mencari data gaji pokok.
"""
        })

    if any(w in query_lower for w in ['cuti', 'izin']):
        education_messages.append({
            "type": "clarification",
            "title": "📋 PERBEDAAN CUTI DAN IZIN",
            "content": """
| Aspek | Cuti | Izin |
| --- | --- | --- |
| Durasi | Panjang: hari-minggu-bulan | Pendek: jam-hari |
| Frekuensi | Terbatas sesuai aturan | Maks 4x/bulan |
| Dokumen | Perlu pengajuan resmi | Cukup alasan sah |

**Jenis Cuti:** Tahunan, Besar, Melahirkan, Sakit, Alasan Penting, Diluar Tanggungan Negara
"""
        })

    return education_messages

def log_rejection(query, reason):
    """Simpan prompt yang ditolak buat audit"""
    try:
        with open(LOG_FILE, "a", encoding="utf-8") as f:
            f.write(f"{datetime.now().isoformat()}|{query}|{reason}\n")
    except Exception as e:
        print(f"[Log] Gagal nulis log: {e}")

# ================== LOADER + SPLITTER ==================
def get_loader(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.pdf':
        try: return UnstructuredPDFLoader(filepath, strategy="hi_res", languages=["ind", "eng"])
        except: return UnstructuredPDFLoader(filepath, strategy="fast", languages=["ind", "eng"])
    elif ext in ['.docx', '.doc']: return Docx2txtLoader(filepath)
    elif ext in ['.xlsx', '.xls', '.csv']: return UnstructuredExcelLoader(filepath, mode="elements")
    elif ext in ['.jpg', '.jpeg', '.png', '.bmp']: return UnstructuredImageLoader(filepath, strategy="ocr_only", languages=["ind", "eng"])
    elif ext == '.txt': return TextLoader(filepath, encoding='utf-8')
    else: raise ValueError(f"Format {ext} belum didukung")

def split_docs(docs):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=600, chunk_overlap=150,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    return splitter.split_documents(docs)

# ================== RAG INDEX ==================
def index_documents(files, progress=gr.Progress()):
    global db, doc_context
    if not files: return "❌ Upload file dulu"
    if os.path.exists(CHROMA_PATH): shutil.rmtree(CHROMA_PATH)

    docs = []
    for f in progress.tqdm(files, desc="Membaca file"):
        try:
            loader = get_loader(f.name)
            docs.extend(loader.load())
        except Exception as e:
            return f"❌ Gagal baca {os.path.basename(f.name)}: {str(e)[:100]}"

    if not docs: return "❌ Tidak ada teks terbaca dari file"

    progress(0.5, desc="Membuat chunk & embedding...")
    chunks = split_docs(docs)
    db = Chroma.from_documents(chunks, embeddings, persist_directory=CHROMA_PATH)

    doc_context = "\n\n".join([d.page_content[:500] for d in docs[:3]])
    sample_text = "\n".join([d.page_content[:200] for d in docs[:2]])

    progress(1.0, desc="Selesai!")
    return f"""✅ Indexing selesai: {len(files)} file → {len(chunks)} chunk. Siap generate!

📄 **Cuplikan dokumen:**
{sample_text}

💡 **Tips:** Sistem akan mengedukasi jika ada istilah yang membingungkan.
"""

# ================== GENERATE + PREVIEW v5.1 PROGRESS BAR ==================
def generate_preview(query, doc_type, mode, progress=gr.Progress()):
    global last_result, doc_context

    if db is None:
        return "❌ Index dokumen dulu di tab Upload", None, gr.update(visible=False), gr.update(visible=False)

    if not query.strip():
        return "❌ Pertanyaan tidak boleh kosong", None, gr.update(visible=False), gr.update(visible=False)

    try:
        # STEP 1: Retrieve
        progress(0.1, desc="Mencari konteks di dokumen...")
        model_name = MODEL_OPTIONS[mode][doc_type]
        retriever = db.as_retriever(search_type="mmr", search_kwargs={"k": 20, "fetch_k": 50})
        docs = retriever.invoke(query)

        if not docs:
            response = "⚠ Data tidak ditemukan dalam dokumen yang diunggah."
            last_result = {"content": response, "type": doc_type}
            return response, None, gr.update(visible=False), gr.update(visible=True)

        context_text = "\n\n".join([d.page_content for d in docs])
        query_lower = query.lower()

        # STEP 2: Edukasi
        progress(0.2, desc="Menganalisis istilah...")
        education_msgs = get_education_context(query, context_text)
        education_html = ""
        if education_msgs:
            education_html = "## 📖 KLARIFIKASI ISTILAH\n\n"
            for msg in education_msgs:
                education_html += f"### {msg['title']}\n{msg['content']}\n\n"

        # STEP 3: Pre-Guard Excel
        progress(0.3, desc="Validasi keamanan data...")
        if doc_type == "excel":
            intent_tabel = any(k in query_lower for k in [
                "tabel", "gaji", "honor", "upah", "pegawai", "nama", "nik",
                "tunjangan", "data", "daftar", "list", "excel", "hitung"
            ])

            table_lines = [l.strip() for l in context_text.split('\n') if '|' in l]
            table_lines = [l for l in table_lines if not re.match(r'^\s*\|[\s\-|:]+\|\s*$', l)]
            table_lines = [l for l in table_lines if l.count('|') >= 2]
            ada_angka_tabel = any(re.search(r'\|\s*\d{1,3}([.,]\d{3})+|\|\s*\d{4,}', l) for l in table_lines)

            if intent_tabel and not (len(table_lines) >= 2 and ada_angka_tabel):
                log_rejection(query, "NO_TABLE_DATA")
                response = f"""{education_html}

⚠ **DITOLAK: Bukan Dokumen Tabel**

Sistem mendeteksi Anda meminta tabel/gaji/data numerik, namun file yang diunggah adalah dokumen narasi peraturan.

📄 **Isi dokumen:** Pasal, ayat, ketentuan, bukan tabel data.

✅ **Yang bisa dilakukan:**
1. Upload file Excel/CSV atau PDF berisi tabel
2. Atau ganti prompt: "Buatkan ringkasan pasal tentang tunjangan" di Tab Word"""
                last_result = {"content": response, "type": doc_type}
                return response, None, gr.update(visible=False), gr.update(visible=True)

        # STEP 4: Panggil LLM
        progress(0.5, desc=f"Memanggil AI {model_name}...")
        system_prompt = f"""Anda adalah asisten RAG profesional dengan pendekatan edukatif.

ATURAN:
1. JAWAB HANYA berdasarkan informasi di KONTEKS.
2. Jika informasi tidak ada, jawab: "Data tidak ditemukan dalam dokumen yang diunggah."
3. Berikan edukasi jika user menggunakan istilah yang keliru.
4. DILARANG KERAS mengarang data apapun.
5. Untuk Excel: Ekstrak tabel dari KONTEKS. JANGAN BUAT DATA BARU.

KONTEKS DOKUMEN:
{context_text}

PERTANYAAN USER: {query}

JAWABAN (edukatif + berdasarkan konteks):"""

        llm = ChatOllama(model=model_name, temperature=0.0, system=system_prompt, num_ctx=8192)
        response = llm.invoke(system_prompt).content.strip()

        # STEP 5: Post-Guard
        progress(0.8, desc="Memvalidasi hasil...")
        banned_names = ["ahmad", "budi", "joko", "siti", "cahya", "deni", "eka", "faisal", "gita", "hari", "indah", "widodo", "ara", "fadil", "awan"]
        if any(f" {name} " in f" {response.lower()} " for name in banned_names):
            log_rejection(query, "HALLUCINATED_NAME")
            response = "⚠ Output diblokir: Sistem mendeteksi halusinasi nama orang."
            last_result = {"content": response, "type": doc_type}
            return response, None, gr.update(visible=False), gr.update(visible=True)

        last_result = {"content": response, "type": doc_type}

        # STEP 6: Parsing Excel
        progress(0.9, desc="Menyiapkan preview...")
        if doc_type == "excel" and "Data tidak ditemukan" not in response and "DITOLAK" not in response:
            try:
                table_lines = [l.strip() for l in response.split('\n') if '|' in l and not re.match(r'^\s*\|[\s\-|:]+\|\s*$', l)]
                if len(table_lines) >= 2:
                    df = pd.read_csv(io.StringIO('\n'.join(table_lines)), sep='|', engine='python', skipinitialspace=True)
                    df = df.dropna(axis=1, how='all')
                    df = df.loc[:, ~df.columns.str.contains('^Unnamed')]
                    df.columns = df.columns.str.strip()
                    progress(1.0, desc="Selesai!")
                    return response, df, gr.update(visible=True), gr.update(visible=True)
            except Exception as e:
                print(f"[Excel] Parse error: {e}")

        progress(1.0, desc="Selesai!")
        return response, None, gr.update(visible=False), gr.update(visible=True)

    except Exception as e:
        print(f"[Generate] Error: {e}")
        return f"❌ Error saat generate: {str(e)[:150]}", None, gr.update(visible=False), gr.update(visible=False)

# ================== EXPORT KE FILE ==================
def export_file():
    if not last_result["content"] or "Data tidak ditemukan" in last_result["content"] or "DITOLAK" in last_result["content"]:
        return None, "❌ Tidak ada konten untuk diexport"

    content = last_result["content"]
    doc_type = last_result["type"]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename_base = f"hasil_{doc_type}_{timestamp}"

    try:
        os.makedirs(OUTPUT_DIR, exist_ok=True)

        if doc_type in ["word", "jurnal", "laporan"]:
            doc = Document()
            doc.add_heading(f'Hasil AI - {doc_type.capitalize()}', 0)
            for line in content.split('\n'):
                if line.startswith('# '): doc.add_heading(line[2:], level=1)
                elif line.startswith('## '): doc.add_heading(line[3:], level=2)
                elif line.strip(): doc.add_paragraph(line)
            path = os.path.join(OUTPUT_DIR, f"{filename_base}.docx")
            doc.save(path)

        elif doc_type == "excel":
            table_lines = [l.strip() for l in content.split('\n') if '|' in l and not re.match(r'^\s*\|[\s\-|:]+\|\s*$', l)]
            if len(table_lines) < 2:
                return None, "❌ Tidak ada tabel untuk diexport"
            df = pd.read_csv(io.StringIO('\n'.join(table_lines)), sep='|', skipinitialspace=True)
            df = df.loc[:, ~df.columns.str.contains('^Unnamed')]
            path = os.path.join(OUTPUT_DIR, f"{filename_base}.xlsx")
            df.to_excel(path, index=False)

        elif doc_type == "ppt":
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = f"Presentasi {doc_type.capitalize()}"
            title_slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%d %B %Y')}"

            sections = re.split(r'\n## ', content)
            for sec in sections[1:7]:
                lines = sec.strip().split('\n')
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = lines[0][:60]
                content_text = '\n'.join([l for l in lines[1:] if l.strip()])[:800]
                slide.placeholders[1].text = content_text

            path = os.path.join(OUTPUT_DIR, f"{filename_base}.pptx")
            prs.save(path)

        return path, f"✅ File berhasil dibuat: {os.path.basename(path)}"
    except Exception as e:
        return None, f"❌ Gagal export: {str(e)[:120]}"

# ================== UI ==================
custom_css = """
.gradio-container {
    max-width: 1400px!important;
    margin: 0 auto!important;
    padding: 0 24px!important;
}

@media (max-width: 768px) {
  .gradio-container {
        max-width: 100%!important;
        padding: 0 12px!important;
    }
    #title {font-size: 1.8em!important;}
}

#title {
    text-align: center;
    background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    font-size: 2.8em;
    font-weight: 800;
    padding: 20px 0 10px 0;
}

#subtitle {
    text-align: center;
    margin-bottom: 24px;
    color: #6b7280;
}

.tabs {flex-wrap: wrap!important;}
.tab-nav button {white-space: nowrap!important;}
.gr-block {
    border-radius: 12px!important;
    border: 1px solid #e5e7eb!important;
}
.gradio-button {
    min-height: 48px!important;
    font-weight: 600!important;
}
.dark.gr-block {border-color: #374151!important;}
.dark #subtitle {color: #9ca3af;}
"""

with gr.Blocks(theme=gr.themes.Soft(primary_hue="violet"), css=custom_css, title="Local AI Trias v5.1") as demo:
    gr.Markdown("# 🤖 Local AI Trias v5.1", elem_id="title")
    gr.Markdown("**Progress Bar + Edukatif + Anti-Halusinasi** | 📚 Sistem Pintar Bedakan Istilah | Ollama Local", elem_id="subtitle")

    with gr.Tab("📁 Upload & Index Dokumen"):
        files = gr.File(file_count="multiple", label="Upload PDF, DOCX, XLSX, JPG, PNG")
        index_btn = gr.Button("🚀 Index Semua File", variant="primary", size="lg")
        index_status = gr.Textbox(label="Status Indexing", interactive=False, max_lines=8)
        index_btn.click(index_documents, inputs=[files], outputs=[index_status])

        gr.Markdown("""
        ### 📚 Sistem Edukasi Istilah Aktif:
        Sistem akan memberikan edukasi jika ada istilah yang membingungkan: gaji vs tunjangan, cuti vs izin.
        **Filosofi:** Lebih baik mengedukasi daripada memblokir!
        """)

    with gr.Tabs() as tabs:
        for doc_type, icon, desc in [
            ("word","📝","Buat Draf Word"),
            ("excel","📊","Ekstrak + Rumus Excel"),
            ("ppt","📽️","Outline PPT"),
            ("jurnal","📚","Ringkasan Jurnal"),
            ("laporan","📋","Laporan Detail")
        ]:
            with gr.Tab(f"{icon} {doc_type.capitalize()}"):
                gr.Markdown(f"### {desc}")

                with gr.Group():
                    with gr.Row():
                        mode = gr.Radio(["Akurat", "Hemat"], value="Akurat", label="Pilih Model", scale=1)
                        query = gr.Textbox(
                            label="Perintah / Pertanyaan",
                            placeholder=f"Contoh: Buatkan {doc_type} tentang...",
                            scale=3,
                            lines=2
                        )

                    with gr.Row():
                        gen_btn = gr.Button("1️⃣ Generate & Preview", variant="primary", scale=2)
                        export_btn = gr.Button("2️⃣ Export ke File", variant="secondary", scale=1)
                        clear_btn = gr.Button("🗑 Clear", variant="secondary", scale=0.5)

                gr.Markdown("#### 📄 Preview Hasil:")
                preview_md = gr.Markdown()
                preview_df = gr.Dataframe(label="Preview Tabel Excel", visible=False, interactive=False, wrap=True)

                with gr.Row():
                    download_file = gr.File(label="⬇ Download File Hasil", visible=False, scale=2)
                    status = gr.Textbox(label="Status", interactive=False, scale=1)

                gen_btn.click(
                    generate_preview,
                    inputs=[query, gr.State(doc_type), mode],
                    outputs=[preview_md, preview_df, preview_df, download_file]
                )
                export_btn.click(export_file, outputs=[download_file, status])
                export_btn.click(lambda: gr.update(visible=True), outputs=[download_file])

                def clear_all():
                    return "", None, gr.update(visible=False), "", gr.update(visible=False)

                clear_btn.click(
                    clear_all,
                    outputs=[query, preview_df, preview_df, preview_md, download_file]
                )

if __name__ == "__main__":
    print(f"🚀 Starting server... Output dir: {OUTPUT_DIR}")
    print(f"📍 Access URL: http://localhost:7860")
    print(f"📚 Mode: Progress Bar + Edukatif + Anti-Halusinasi")
    demo.launch(
        server_name="127.0.0.1",
        server_port=7860,
        allowed_paths=[OUTPUT_DIR],
        show_error=True
    )