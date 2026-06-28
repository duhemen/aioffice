import gradio as gr
import ollama
from docx import Document
import pandas as pd
from pptx import Presentation
import os
from datetime import datetime
import fitz # PyMuPDF
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from io import StringIO

# ================== KONFIGURASI ==================
MODEL_INSTRUCT = 'qwen2.5:14b-instruct-q4_K_M'
OUTPUT_DIR = "hasil_generate"
DB_DIR = "chroma_db"
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(DB_DIR, exist_ok=True)

# Text splitter GLOBAL - yang bener, dipake di indexing
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=2000, # Lebih besar biar tabel muat 1 chunk
    chunk_overlap=400, # Overlap gede biar nggak kepotong
    separators=["\n\n", "\n", "Pasal", "BAB", "LAMPIRAN", " "]
)

# 1. SETUP EMBEDDING - pake CPU biar ringan, ganti 'cuda' kalau mau
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    model_kwargs={'device': 'cpu'}
)

# 2. FUNGSI BACA FILE
def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext == '.pdf':
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += f"\n--- HALAMAN {page.number + 1} ---\n" + page.get_text()
        elif ext == '.docx':
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == '.xlsx':
            df = pd.read_excel(file_path, sheet_name=None)
            for name, sheet in df.items():
                text += f"Sheet: {name}\n{sheet.to_markdown(index=False)}\n\n"
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text += f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error baca file: {str(e)}"

# 3. INDEXING
def index_documents(files):
    if not files:
        return "❌ Upload file dulu"
    all_texts = []
    metadatas = []
    for file in files:
        text = extract_text_from_file(file.name)
        all_texts.append(text)
        metadatas.append({"source": os.path.basename(file.name)})

    docs = text_splitter.create_documents(all_texts, metadatas=metadatas)

    # Hapus DB lama
    if os.path.exists(DB_DIR):
        import shutil
        shutil.rmtree(DB_DIR)

    vectordb = Chroma.from_documents(documents=docs, embedding=embeddings, persist_directory=DB_DIR)
    vectordb.persist()
    return f"✅ Berhasil index {len(files)} file, {len(docs)} chunks. Siap buat RAG!"

# 4. RETRIEVAL
def get_context(query):
    try:
        vectordb = Chroma(persist_directory=DB_DIR, embedding_function=embeddings)
        docs = vectordb.similarity_search(query, k=8) # Ambil 8 chunk biar lengkap
        context = "\n\n".join([f"[Sumber: {d.metadata['source']}]\n{d.page_content}" for d in docs])
        return context
    except:
        return ""

# 5. OLLAMA OPTIONS - ANTI BLUP BLUP
OLLAMA_OPTIONS = {
    "num_ctx": 8192,
    "num_gpu": 999, # Paksa pake RTX 4060 8GB
    "temperature": 0.0, # 0 = anti ngarang
    "top_p": 0.9
}

# 6. GENERATE WORD - PROMPT ANTI HALUSINASI
def generate_word(judul, perintah):
    if not judul.strip():
        return None, "❌ Judul laporan wajib diisi"
    try:
        context = get_context(perintah)

        # PROMPT GALAK ANTI NGARANG
        prompt = f"""Kamu adalah mesin ekstrak data. DILARANG MENGARANG.

KONTEKS DOKUMEN:
---
{context}
---

TUGAS: Buatkan dokumen formal judul '{judul}'.
INSTRUKSI: {perintah}

ATURAN KERAS:
1. Jawab HANYA dari konteks di atas. Copy paste teks/tabel apa adanya.
2. Jika data tidak ada di konteks, tulis persis: "Data tidak ditemukan di dokumen".
3. DILARANG membuat nama, angka, atau contoh fiktif seperti "Ibu Ara".
4. Untuk tabel, pertahankan format markdown. Untuk angka Rp, copy persis.
5. Gunakan Bahasa Indonesia baku formal.

OUTPUT:"""

        response = ollama.chat(model=MODEL_INSTRUCT, messages=[{'role': 'user', 'content': prompt}], options=OLLAMA_OPTIONS)
        hasil_ai = response['message']['content']

        doc = Document()
        doc.add_heading(judul, 0)
        for baris in hasil_ai.split('\n'):
            if baris.strip():
                doc.add_paragraph(baris)

        filename = f"{OUTPUT_DIR}/{judul.replace(' ', '_')}_{datetime.now().strftime('%H%M%S')}.docx"
        doc.save(filename)
        return filename, f"✅ Berhasil! RAG: {'Aktif ('+str(len(context))+' char)' if context else 'Non-aktif'}"
    except Exception as e:
        return None, f"❌ Error: {str(e)}"

def generate_excel(judul, perintah):
    try:
        context = get_context(perintah)
        prompt = f"""Berdasarkan konteks berikut, ekstrak HANYA tabel yang ada.
---
{context}
---
Instruksi: {perintah}
Output HANYA format CSV dengan header di baris pertama. Jika tabel tidak ada, output: "kolom1,kolom2\nData tidak ditemukan,Data tidak ditemukan"
Jangan beri penjelasan."""

        response = ollama.chat(model=MODEL_INSTRUCT, messages=[{'role': 'user', 'content': prompt}], options=OLLAMA_OPTIONS)
        csv_data = response['message']['content'].strip().replace('```csv','').replace('```','')

        df = pd.read_csv(StringIO(csv_data))
        filename = f"{OUTPUT_DIR}/{judul.replace(' ', '_')}_{datetime.now().strftime('%H%M%S')}.xlsx"
        df.to_excel(filename, index=False)
        return filename, f"✅ Berhasil! RAG: {'Aktif' if context else 'Non-aktif'}"
    except Exception as e:
        return None, f"❌ Error: {str(e)}. Coba cek prompt."

def generate_ppt(judul, perintah):
    try:
        context = get_context(perintah)
        prompt = f"""Berdasarkan konteks berikut, buat outline PPT.
---
{context}
---
Judul: {judul}. Instruksi: {perintah}.
Format: Judul Slide 1:... | Poin 1 | Poin 2 || Judul Slide 2:...
Hanya output format itu."""

        response = ollama.chat(model=MODEL_INSTRUCT, messages=[{'role': 'user', 'content': prompt}], options=OLLAMA_OPTIONS)
        hasil_ai = response['message']['content']

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = judul
        slide.placeholders[1].text = "aioffice Pro RAG"

        slides_data = hasil_ai.split('||')
        for slide_text in slides_data:
            if 'Judul Slide' in slide_text:
                parts = [p.strip() for p in slide_text.split('|') if p.strip()]
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = parts[0].replace('Judul Slide', '').replace(':', '').strip()
                body = slide.shapes.placeholders[1].text_frame
                body.clear()
                for poin in parts[1:]:
                    p = body.add_paragraph()
                    p.text = poin
                    p.level = 0

        filename = f"{OUTPUT_DIR}/{judul.replace(' ', '_')}_{datetime.now().strftime('%H%M%S')}.pptx"
        prs.save(filename)
        return filename, f"✅ Berhasil! RAG: {'Aktif' if context else 'Non-aktif'}"
    except Exception as e:
        return None, f"❌ Error: {str(e)}"

# CSS
custom_css = """
#main_header {text-align: center; background: linear-gradient(90deg, #667eea 0%, #764ba2 100%); padding: 25px; border-radius: 15px; margin-bottom: 20px; width: 100%;}
#main_header h1 {color: white; margin: 0; font-size: 2.2em;}
#main_header p {color: #e0e0e0; margin: 8px 0 0 0; font-size: 1.1em;}
.gradio-container {max-width: 100%!important; margin: 0!important; padding: 20px!important;}
"""

with gr.Blocks(theme=gr.themes.Soft(primary_hue="violet"), css=custom_css, title="aioffice Pro", fill_width=True) as demo:
    with gr.Row(elem_id="main_header"):
        gr.HTML("<h1>🚀 aioffice Pro v2.0 RAG Anti-Halu</h1><p>Qwen2.5:14b + RTX 4060 Optimized</p>")

    with gr.Tabs():
        with gr.TabItem("📚 RAG - Upload"):
            file_upload = gr.File(label="Upload PDF/DOCX/XLSX/PPTX", file_count="multiple", file_types=[".pdf",".docx",".xlsx",".pptx"])
            btn_index = gr.Button("1. Index File", variant="primary")
            status_index = gr.Textbox(label="Status")
            btn_index.click(index_documents, inputs=file_upload, outputs=status_index)

        with gr.TabItem("📄 Generate Word"):
            judul_word = gr.Textbox(label="Judul", placeholder="Laporan Lengkap PMPU 9 2026")
            perintah_word = gr.Textbox(label="Perintah", lines=6, placeholder="Ekstrak semua tabel contoh perhitungan Bapak Fadil dari Lampiran A. Jangan diubah.")
            btn_word = gr.Button("Generate", variant="primary")
            file_word = gr.File(label="Download")
            status_word = gr.Textbox(label="Status")
            btn_word.click(generate_word, [judul_word, perintah_word], [file_word, status_word])

        with gr.TabItem("📊 Generate Excel"):
            judul_excel = gr.Textbox(label="Judul")
            perintah_excel = gr.Textbox(label="Perintah", lines=4)
            btn_excel = gr.Button("Generate", variant="primary")
            file_excel = gr.File(label="Download")
            status_excel = gr.Textbox(label="Status")
            btn_excel.click(generate_excel, [judul_excel, perintah_excel], [file_excel, status_excel])

        with gr.TabItem("📽 Generate PPT"):
            judul_ppt = gr.Textbox(label="Judul")
            perintah_ppt = gr.Textbox(label="Perintah", lines=4)
            btn_ppt = gr.Button("Generate", variant="primary")
            file_ppt = gr.File(label="Download")
            status_ppt = gr.Textbox(label="Status")
            btn_ppt.click(generate_ppt, [judul_ppt, perintah_ppt], [file_ppt, status_ppt])

if __name__ == "__main__":
    # PENTING: Set environment sebelum launch biar GPU kepake
    os.environ["OLLAMA_NUM_GPU"] = "999"
    demo.launch(server_name="127.0.0.1", server_port=7860)