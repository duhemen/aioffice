import gradio as gr
import ollama
import os
import yaml
import pandas as pd
from docx import Document
from pptx import Presentation
from datetime import datetime
import fitz  # PyMuPDF
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_huggingface import HuggingFaceEmbeddings  # <-- FIX warning
from io import StringIO

# ========== KONFIG TRIAS ==========
CONFIG = yaml.safe_load(open("models.yaml", encoding="utf-8"))
MODEL_WORD = CONFIG['models']['word']['name']   # qwen2.5:14b
MODEL_EXCEL = CONFIG['models']['excel']['name'] # deepseek-coder:6.7b
MODEL_PPT = CONFIG['models']['ppt']['name']     # mistral:7b

# LOGO - auto detect jpg/png
LOGO = "assets/logo.jpeg" if os.path.exists("assets/logo.jpeg") else "assets/logo.png"

OUTPUT_DIR = "hasil_generate"
DB_DIR = "chroma_db"
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(DB_DIR, exist_ok=True)

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=2000, 
    chunk_overlap=400,
    separators=["\n\n", "\n", "Pasal", "BAB", "LAMPIRAN", " "]
)

embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    model_kwargs={'device': 'cpu'}
)

OLLAMA_OPTS = {"num_ctx": 8192, "num_gpu": 999, "temperature": 0.0, "top_p": 0.9}
def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext=='.pdf':
        with fitz.open(file_path) as doc: return "\n".join([f"---HAL {p.number+1}---\n{p.get_text()}" for p in doc])
    if ext=='.docx': return "\n".join([p.text for p in Document(file_path).paragraphs])
    if ext=='.xlsx':
        dfs = pd.read_excel(file_path, sheet_name=None)
        return "\n\n".join([f"Sheet:{n}\n{df.to_markdown(index=False)}" for n,df in dfs.items()])
    if ext=='.pptx':
        prs = Presentation(file_path); return "\n".join([s.shapes.title.text if s.shapes.title else "" for s in prs.slides])
    return ""

def index_docs(files):
    if not files: return "❌ Upload dulu"
    import shutil; shutil.rmtree(DB_DIR, ignore_errors=True)
    texts, metas = [], []
    for f in files: texts.append(extract_text(f.name)); metas.append({"source":os.path.basename(f.name)})
    docs = text_splitter.create_documents(texts, metadatas=metas)
    Chroma.from_documents(docs, embeddings, persist_directory=DB_DIR).persist()
    return f"✅ TRIAS RAG aktif! {len(files)} file, {len(docs)} chunks"

def get_ctx(q):
    try: db=Chroma(persist_directory=DB_DIR, embedding_function=embeddings); return "\n\n".join([d.page_content for d in db.similarity_search(q,k=8)])
    except: return ""

def call(model, prompt): return ollama.chat(model=model, messages=[{"role":"user","content":prompt}], options=OLLAMA_OPTS)['message']['content']

def gen_word(judul, perintah):
    ctx=get_ctx(perintah)
    p=f"""Kamu mesin ekstrak. DILARANG NGARANG.
KONTEKS:
{ctx}
TUGAS: {judul}. INSTRUKSI: {perintah}
ATURAN: 1)Copy persis dari konteks 2)Jika tidak ada tulis "Data tidak ditemukan" 3)Bahasa Indonesia formal
OUTPUT:"""
    out = call(MODEL_WORD, p)
    doc=Document(); doc.add_heading(judul,0)
    for ln in out.split("\n"): doc.add_paragraph(ln)
    fn=f"{OUTPUT_DIR}/{judul[:30].replace(' ','_')}_{datetime.now().strftime('%H%M')}.docx"; doc.save(fn)
    return fn, f"✅ Qwen selesai | RAG:{'ON' if ctx else 'OFF'}"

def gen_excel(judul, perintah):
    ctx=get_ctx(perintah)
    p=f"Konteks:{ctx}\nEkstrak tabel jadi CSV header pertama. Perintah:{perintah}\nHanya CSV:"
    csv = call(MODEL_EXCEL, p).replace("```csv","").replace("```","")
    df=pd.read_csv(StringIO(csv)); fn=f"{OUTPUT_DIR}/{judul[:30]}.xlsx"; df.to_excel(fn,index=False)
    return fn, f"✅ DeepSeek selesai | RAG:{'ON' if ctx else 'OFF'}"

def gen_ppt(judul, perintah):
    ctx=get_ctx(perintah)
    p=f"Konteks:{ctx}\nBuat outline PPT judul '{judul}'. Instruksi:{perintah}\nFormat: Judul1|poin|poin||Judul2|..."
    out=call(MODEL_PPT,p)
    prs=Presentation(); s=prs.slides.add_slide(prs.slide_layouts[0]); s.shapes.title.text=judul; s.placeholders[1].text="TILAIOFFICE TRIAS"
    for sd in out.split("||"):
        parts=[x.strip() for x in sd.split("|") if x.strip()]
        if parts: sl=prs.slides.add_slide(prs.slide_layouts[1]); sl.shapes.title.text=parts[0]; tf=sl.placeholders[1].text_frame; tf.clear()
        for pt in parts[1:]: p=tf.add_paragraph(); p.text=pt
    fn=f"{OUTPUT_DIR}/{judul[:30]}.pptx"; prs.save(fn)
    return fn, f"✅ Mistral selesai | RAG:{'ON' if ctx else 'OFF'}"

css = """.header{background:linear-gradient(90deg,#667eea,#764ba2);padding:25px;border-radius:15px;color:white;text-align:center;margin-bottom:20px}"""

with gr.Blocks(title="TILAIOFFICE TRIAS") as demo:
    # HEADER + LOGO (logo.jpeg)
    gr.HTML(f"""
    <div class="header">
        <img src="/gradio_api/file={LOGO}" style="height:60px;vertical-align:middle;margin-right:12px" onerror="this.style.display='none'">
        <div style="display:inline-block;vertical-align:middle;text-align:left">
            <h1 style="margin:0;font-size:26px">🚀 TILAIOFFICE Pro v3 TRIAS</h1>
            <p style="margin:3px 0 0 0">Qwen + DeepSeek + Mistral | RAG Anti-Halu | RTX 4060</p>
            <p style="margin:2px 0 0 0;font-size:0.85em;opacity:0.85">Created by Your Self</p>
        </div>
    </div>
    """)
    
    with gr.Tabs():
        # TAB 1 - UPLOAD
        with gr.TabItem("📚 1. Upload RAG"):
            up = gr.File(file_count="multiple", file_types=[".pdf",".docx",".xlsx",".pptx"], label="Upload dokumen sumber")
            b = gr.Button("Index ke TRIAS", variant="primary")
            st = gr.Textbox(label="Status")
            b.click(index_docs, up, st)
        
        # TAB 2 - WORD
        with gr.TabItem("📄 2. Word (Qwen)"):
            jw = gr.Textbox(label="Judul")
            pw = gr.Textbox(label="Perintah", lines=5)
            bw = gr.Button("Generate", variant="primary")
            fw = gr.File(label="Download")
            sw = gr.Textbox(label="Status")
            bw.click(gen_word, [jw, pw], [fw, sw])
        
        # TAB 3 - EXCEL
        with gr.TabItem("📊 3. Excel (DeepSeek)"):
            je = gr.Textbox(label="Judul")
            pe = gr.Textbox(label="Perintah", lines=4)
            be = gr.Button("Generate", variant="primary")
            fe = gr.File(label="Download")
            se = gr.Textbox(label="Status")
            be.click(gen_excel, [je, pe], [fe, se])
        
        # TAB 4 - PPT
        with gr.TabItem("🎬 4. PPT (Mistral)"):
            jp = gr.Textbox(label="Judul")
            pp = gr.Textbox(label="Perintah", lines=4)
            bp = gr.Button("Generate", variant="primary")
            fp = gr.File(label="Download")
            sp = gr.Textbox(label="Status")
            bp.click(gen_ppt, [jp, pp], [fp, sp])
    
    # FOOTER - yang hilang
    gr.HTML("""
    <div style='text-align:center;margin-top:40px;padding:20px;color:#666;font-size:0.9em'>
    Dibuat dengan ❤️ <b>Ingat LUCA</b> | <a href='#' style='color:#764ba2;text-decoration:none'>GitHub Repo</a> | aioffice v3.0<br>
    100% Offline • Data Aman di Laptop Sendiri • Powered by Your Self
    </div>
    """)

if __name__=="__main__":
    os.environ["OLLAMA_NUM_GPU"]="999"
    demo.launch(
    server_name="127.0.0.1",
    server_port=7860,
    allowed_paths=[os.path.abspath("assets")],
    css=css,
    theme=gr.themes.Soft(primary_hue="violet")
)